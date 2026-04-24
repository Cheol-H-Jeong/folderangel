"""Parsers for modern Office / ODT formats."""
from __future__ import annotations

import logging
import zipfile
from pathlib import Path
from xml.etree import ElementTree as ET

# (NS constant defined further down next to the PPTX helpers.)

log = logging.getLogger(__name__)


def _cap(chunks: list[str], max_chars: int) -> str:
    out = "\n".join(c for c in chunks if c)
    return out[:max_chars]


def parse_docx(path: Path, max_chars: int) -> str:
    try:
        from docx import Document  # type: ignore
    except ImportError:
        log.warning("python-docx not installed; skipping %s", path)
        return ""
    try:
        doc = Document(str(path))
    except Exception as exc:
        log.warning("docx open failed %s: %s", path, exc)
        return ""
    chunks: list[str] = []
    total = 0
    for p in doc.paragraphs:
        txt = (p.text or "").strip()
        if not txt:
            continue
        chunks.append(txt)
        total += len(txt)
        if total >= max_chars:
            break
    # Include the first table as a fallback if paragraphs were empty
    if total < 40:
        for tbl in doc.tables[:2]:
            for row in tbl.rows:
                row_txt = " | ".join((c.text or "").strip() for c in row.cells)
                if row_txt.strip():
                    chunks.append(row_txt)
                    total += len(row_txt)
                    if total >= max_chars:
                        break
            if total >= max_chars:
                break
    return _cap(chunks, max_chars)


_PPTX_A_NS = "{http://schemas.openxmlformats.org/drawingml/2006/main}"


def _parse_pptx_via_xml(path: Path, max_chars: int) -> str:
    """Fallback PPTX text extractor that reads slide XML directly.

    Used when python-pptx's high-level API trips over malformed relationship
    parts (seen with some Korean templates: ``'list' object has no
    attribute 'rId'``).  We only need ``<a:t>`` text runs to feed the LLM.
    """
    chunks: list[str] = []
    total = 0
    try:
        with zipfile.ZipFile(path) as zf:
            slide_names = sorted(
                n for n in zf.namelist()
                if n.startswith("ppt/slides/slide") and n.endswith(".xml")
            )
            for name in slide_names[:30]:
                try:
                    with zf.open(name) as f:
                        xml = f.read()
                except Exception:
                    continue
                try:
                    root = ET.fromstring(xml)
                except ET.ParseError:
                    continue
                for elem in root.iter(f"{_PPTX_A_NS}t"):
                    txt = (elem.text or "").strip()
                    if not txt:
                        continue
                    chunks.append(txt)
                    total += len(txt)
                    if total >= max_chars:
                        return _cap(chunks, max_chars)
    except (zipfile.BadZipFile, KeyError) as exc:
        log.warning("pptx xml fallback failed %s: %s", path, exc)
    return _cap(chunks, max_chars)


def parse_pptx(path: Path, max_chars: int) -> str:
    try:
        from pptx import Presentation  # type: ignore
    except ImportError:
        return _parse_pptx_via_xml(path, max_chars)
    try:
        pres = Presentation(str(path))
    except Exception as exc:
        log.debug("pptx open failed %s: %s — falling back to xml", path, exc)
        return _parse_pptx_via_xml(path, max_chars)
    chunks: list[str] = []
    total = 0
    try:
        for slide in pres.slides[:20]:
            for shape in slide.shapes:
                if not getattr(shape, "has_text_frame", False):
                    continue
                try:
                    paragraphs = shape.text_frame.paragraphs
                except Exception:
                    continue
                for para in paragraphs:
                    try:
                        runs = para.runs
                    except Exception:
                        continue
                    for run in runs:
                        try:
                            txt = (run.text or "").strip()
                        except Exception:
                            continue
                        if not txt:
                            continue
                        chunks.append(txt)
                        total += len(txt)
                        if total >= max_chars:
                            return _cap(chunks, max_chars)
    except Exception as exc:
        # Some malformed decks raise deep inside python-pptx (e.g. relationship
        # objects returning lists instead of part refs).  Fall back to the raw
        # XML extractor so we still get usable text.
        log.debug("pptx high-level walk failed %s: %s — falling back to xml", path, exc)
        if not chunks:
            return _parse_pptx_via_xml(path, max_chars)
    if total < 40:
        # If high-level walk produced nothing meaningful, try the XML fallback.
        xml_text = _parse_pptx_via_xml(path, max_chars)
        if len(xml_text) > total:
            return xml_text
    return _cap(chunks, max_chars)


def parse_xlsx(path: Path, max_chars: int) -> str:
    try:
        from openpyxl import load_workbook  # type: ignore
    except ImportError:
        log.warning("openpyxl not installed; skipping %s", path)
        return ""
    try:
        wb = load_workbook(str(path), read_only=True, data_only=True)
    except Exception as exc:
        log.warning("xlsx open failed %s: %s", path, exc)
        return ""
    chunks: list[str] = []
    total = 0
    for name in wb.sheetnames[:4]:
        chunks.append(f"# Sheet: {name}")
        ws = wb[name]
        for row in ws.iter_rows(max_row=15, values_only=True):
            cells = [str(c) for c in row if c is not None]
            if not cells:
                continue
            row_txt = " | ".join(cells)
            chunks.append(row_txt)
            total += len(row_txt)
            if total >= max_chars:
                return _cap(chunks, max_chars)
    return _cap(chunks, max_chars)


def parse_odt(path: Path, max_chars: int) -> str:
    try:
        with zipfile.ZipFile(path) as z:
            with z.open("content.xml") as f:
                xml = f.read()
    except (zipfile.BadZipFile, KeyError) as exc:
        log.warning("odt open failed %s: %s", path, exc)
        return ""
    try:
        root = ET.fromstring(xml)
    except ET.ParseError as exc:
        log.warning("odt xml parse failed %s: %s", path, exc)
        return ""
    texts: list[str] = []
    total = 0
    for elem in root.iter():
        if elem.text and elem.text.strip():
            texts.append(elem.text.strip())
            total += len(elem.text)
            if total >= max_chars:
                break
    return _cap(texts, max_chars)
